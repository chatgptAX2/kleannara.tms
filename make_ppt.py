"""
TMS 개발 진행현황 보고자료 PPT 생성 (2슬라이드)
슬라이드1: TMS 시스템 소개 + 도입효과
슬라이드2: 개발 일정 (간트차트 중심)
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── 색상 팔레트 ────────────────────────────────────────────────
C_NAVY   = RGBColor(0x0D, 0x2B, 0x55)
C_BLUE   = RGBColor(0x16, 0x4E, 0xA0)
C_ACCENT = RGBColor(0x00, 0xA8, 0xE8)
C_GREEN  = RGBColor(0x0F, 0x7B, 0x6C)
C_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
C_PURPLE = RGBColor(0x6D, 0x28, 0xD9)
C_GRAY   = RGBColor(0x4B, 0x55, 0x63)
C_LGRAY  = RGBColor(0xF1, 0xF5, 0xF9)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK   = RGBColor(0x1E, 0x29, 0x3B)
C_SILVER = RGBColor(0xE2, 0xE8, 0xF0)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── 공통 헬퍼 ─────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def tb(slide, text, x, y, w, h, size=11, bold=False, color=C_DARK,
       align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame; tf.word_wrap = wrap
    p   = tf.paragraphs[0]; p.alignment = align
    r   = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return txb

def header_bar(slide, title, subtitle, page_n):
    """상단 공통 헤더"""
    W = SLIDE_W
    rect(slide, 0, 0, W, SLIDE_H, RGBColor(0xF6, 0xF8, 0xFB))   # 전체 배경
    rect(slide, 0, 0, W, Inches(1.08), C_NAVY)                    # 헤더 바
    rect(slide, 0, 0, Inches(0.07), SLIDE_H, C_ACCENT)            # 좌측 세로선
    tb(slide, "kleannara TMS", Inches(0.22), Inches(0.07),
       Inches(3), Inches(0.32), size=10, bold=True,
       color=RGBColor(0x90, 0xCA, 0xF9))
    tb(slide, title, Inches(0.22), Inches(0.28), Inches(10), Inches(0.52),
       size=21, bold=True, color=C_WHITE)
    tb(slide, subtitle, Inches(0.22), Inches(0.78), Inches(7), Inches(0.26),
       size=9.5, color=RGBColor(0xBF, 0xDB, 0xFF))
    tb(slide, page_n, Inches(12.2), Inches(0.80), Inches(1.0), Inches(0.24),
       size=9.5, color=RGBColor(0xBF, 0xDB, 0xFF), align=PP_ALIGN.RIGHT)

def section_card(slide, x, y, w, h, hdr_text, hdr_color=C_BLUE):
    """섹션 카드 (흰 배경 + 컬러 헤더 바)"""
    rect(slide, x, y, w, h, C_WHITE)
    rect(slide, x, y, w, Inches(0.36), hdr_color)
    tb(slide, hdr_text, x + Inches(0.12), y + Inches(0.05),
       w - Inches(0.2), Inches(0.28), size=11, bold=True, color=C_WHITE)


# ════════════════════════════════════════════════════════════════
#  슬라이드 1: TMS 시스템 소개 + 도입효과
# ════════════════════════════════════════════════════════════════
def build_slide1(prs):
    slide = blank_slide(prs)
    header_bar(slide,
               "TMS (Transport Management System) 시스템 소개",
               "개발팀  |  보고일: 2026년 5월",
               "1 / 2")

    # ── 좌측 패널 ── TMS란? ─────────────────────────────────
    px, py = Inches(0.18), Inches(1.18)
    pw, ph = Inches(5.55), Inches(5.95)
    section_card(slide, px, py, pw, ph, "📌  TMS란 무엇인가?", C_NAVY)

    # 정의
    cy = py + Inches(0.48)
    rect(slide, px + Inches(0.14), cy, pw - Inches(0.28), Inches(0.78),
         RGBColor(0xEF, 0xF6, 0xFF))
    tb(slide, "TMS (Transport Management System)",
       px + Inches(0.22), cy + Inches(0.06), pw - Inches(0.44), Inches(0.26),
       size=11, bold=True, color=C_BLUE)
    tb(slide,
       "물류·유통 기업의 출하 계획부터 배차 지시, 운송 실행,\n"
       "운임 정산까지 전 과정을 통합 관리하는 운송관리 시스템",
       px + Inches(0.22), cy + Inches(0.30), pw - Inches(0.44), Inches(0.44),
       size=9.5, color=C_DARK, wrap=True)
    cy += Inches(0.90)

    # 시스템 위치도 (ERP ─▶ TMS ─▶ 현장)
    bar_items = [
        ("ERP\n(출고지시)",  C_BLUE,   RGBColor(0xDB, 0xEA, 0xFF)),
        ("TMS\n운송관리",    C_NAVY,   RGBColor(0xD0, 0xD9, 0xF0)),
        ("WMS\n창고관리",    C_ACCENT, RGBColor(0xD0, 0xF0, 0xFF)),
        ("현장·차량\n(운송)", C_GREEN,  RGBColor(0xD1, 0xFA, 0xEA)),
    ]
    bw = (pw - Inches(0.42)) / len(bar_items)
    for bi, (label, fc, bg) in enumerate(bar_items):
        bx2 = px + Inches(0.14) + bi * (bw + Inches(0.04))
        rect(slide, bx2, cy, bw, Inches(0.54), bg)
        rect(slide, bx2, cy, bw, Inches(0.06), fc)
        tb(slide, label, bx2, cy + Inches(0.08), bw, Inches(0.44),
           size=8.5, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
        if bi < len(bar_items) - 1:
            tb(slide, "▶", bx2 + bw, cy + Inches(0.16), Inches(0.1), Inches(0.28),
               size=9, color=C_GRAY, align=PP_ALIGN.CENTER)
    cy += Inches(0.66)

    # 핵심 기능 리스트
    features = [
        ("🗓️", "배차 계획 자동화",
         "출고 데이터를 기반으로 차량 용량·경로·배송 우선순위를\n"
         "자동 계산하여 최적 배차안 생성"),
        ("🚛", "차량·운전자 관리",
         "차량 제원(톤수·적재중량·파렛트수) 및 운전자 정보를\n"
         "통합 관리하고 배차에 자동 반영"),
        ("📍", "납품처·경로 관리",
         "납품처별 배송구역·우선순위·경유지 정보를 등록하고\n"
         "운송경로별 비용 기준으로 운임 자동 산출"),
        ("📊", "실적·비용 분석",
         "배차 실적, 노선별 운임, 차량 가동률 등 KPI를\n"
         "대시보드에서 실시간 모니터링"),
    ]
    for icon, title, desc in features:
        rect(slide, px + Inches(0.14), cy, pw - Inches(0.28), Inches(0.72),
             RGBColor(0xF8, 0xFA, 0xFC))
        # 아이콘 박스
        rect(slide, px + Inches(0.14), cy, Inches(0.44), Inches(0.72), C_SILVER)
        tb(slide, icon, px + Inches(0.14), cy + Inches(0.16),
           Inches(0.44), Inches(0.38), size=14, align=PP_ALIGN.CENTER)
        tb(slide, title, px + Inches(0.64), cy + Inches(0.05),
           pw - Inches(0.82), Inches(0.24), size=9.5, bold=True, color=C_BLUE)
        tb(slide, desc, px + Inches(0.64), cy + Inches(0.26),
           pw - Inches(0.82), Inches(0.42), size=8.5, color=C_GRAY, wrap=True)
        cy += Inches(0.78)

    # ── 우측 패널 ── 도입효과 ───────────────────────────────
    rx, ry = Inches(5.90), Inches(1.18)
    rw, rh = Inches(7.25), Inches(5.95)
    section_card(slide, rx, ry, rw, rh, "✨  TMS 도입 기대효과", C_BLUE)

    effects = [
        {
            "no": "01",
            "title": "배차 업무 효율화",
            "color": C_BLUE,
            "bg":    RGBColor(0xEB, 0xF3, 0xFF),
            "icon":  "⚡",
            "points": [
                "수작업 배차 → 자동배차 전환으로 담당자 업무 부하 대폭 감소",
                "배차 처리 시간 단축 (목표: 기존 대비 70% 이상 절감)",
                "배차 오류·누락 사전 차단 및 배차 품질 표준화",
            ]
        },
        {
            "no": "02",
            "title": "운송비 절감",
            "color": C_GREEN,
            "bg":    RGBColor(0xE6, 0xF7, 0xF4),
            "icon":  "💰",
            "points": [
                "적재율 최적화(혼적·분할 배차)로 공차율 및 운송비 절감",
                "경로 최적화로 불필요한 이동 거리·시간 단축",
                "운송경로별 비용 기준 자동 정산으로 운임 투명성 확보",
            ]
        },
        {
            "no": "03",
            "title": "고객 서비스 향상",
            "color": C_ACCENT,
            "bg":    RGBColor(0xE0, 0xF5, 0xFF),
            "icon":  "🎯",
            "points": [
                "납품처별 배송 우선순위 반영으로 납기 준수율 향상",
                "배송 현황 실시간 조회 및 이상 발생 시 즉각 대응",
                "정확한 배송 예정 시간(ETA) 제공으로 고객 신뢰도 제고",
            ]
        },
        {
            "no": "04",
            "title": "경영 의사결정 지원",
            "color": C_PURPLE,
            "bg":    RGBColor(0xF3, 0xEE, 0xFF),
            "icon":  "📈",
            "points": [
                "배차 실적·운임·차량 가동률 KPI 대시보드 실시간 제공",
                "노선별·차급별 비용 분석으로 운송 전략 수립 지원",
                "ERP·WMS 데이터 연동으로 물류 전체 가시성 확보",
            ]
        },
    ]

    ey = ry + Inches(0.48)
    ew = rw - Inches(0.28)
    for eff in effects:
        eh = Inches(1.24)
        rect(slide, rx + Inches(0.14), ey, ew, eh, eff["bg"])
        # 번호 + 제목 바
        rect(slide, rx + Inches(0.14), ey, ew, Inches(0.30), eff["color"])
        tb(slide, f"{eff['no']}  {eff['title']}",
           rx + Inches(0.22), ey + Inches(0.03), ew - Inches(0.16), Inches(0.25),
           size=10, bold=True, color=C_WHITE)
        # 아이콘
        tb(slide, eff["icon"], rx + Inches(0.14), ey + Inches(0.31),
           Inches(0.36), Inches(0.56), size=18, align=PP_ALIGN.CENTER)
        # 포인트
        for pi, pt in enumerate(eff["points"]):
            tb(slide, f"•  {pt}",
               rx + Inches(0.54), ey + Inches(0.30) + pi * Inches(0.30),
               ew - Inches(0.56), Inches(0.28),
               size=8.5, color=C_DARK, wrap=True)
        ey += eh + Inches(0.1)

    return slide


# ════════════════════════════════════════════════════════════════
#  슬라이드 2: 개발 일정 (간트차트 중심)
# ════════════════════════════════════════════════════════════════
def build_slide2(prs):
    slide = blank_slide(prs)
    header_bar(slide,
               "TMS 개발 일정 및 단계별 추진 계획",
               "개발팀  |  기준일: 2026년 5월",
               "2 / 2")

    W, H = SLIDE_W, SLIDE_H

    # ── 개요 배너 (2단계 요약) ─────────────────────────────
    bx, by = Inches(0.18), Inches(1.18)
    bw, bh = W - Inches(0.36), Inches(0.78)
    rect(slide, bx, by, bw, bh, RGBColor(0xF0, 0xF4, 0xFF))
    rect(slide, bx, by, Inches(0.06), bh, C_NAVY)

    phase_summary = [
        ("1단계", "배차 효율화 개발",        "자동배차 기능을 통한 배차 효율화 향상", "2026. 04 ~ 06",  C_GREEN),
        ("2단계", "차량관제 및 운송비 관리 개발", "실시간 차량관제·운임 정산·실적 분석", "2026. 07 ~ 12",  C_ORANGE),
    ]
    sw = (bw - Inches(0.2)) / 2
    for si, (phase, title, desc, period, pc) in enumerate(phase_summary):
        sx = bx + Inches(0.1) + si * sw
        rect(slide, sx + Inches(0.06), by + Inches(0.1), Inches(0.52), Inches(0.55), pc)
        tb(slide, phase, sx + Inches(0.06), by + Inches(0.1),
           Inches(0.52), Inches(0.55), size=9, bold=True,
           color=C_WHITE, align=PP_ALIGN.CENTER)
        tb(slide, title, sx + Inches(0.64), by + Inches(0.10),
           sw - Inches(0.72), Inches(0.26), size=10.5, bold=True, color=C_DARK)
        tb(slide, desc, sx + Inches(0.64), by + Inches(0.32),
           sw - Inches(0.72), Inches(0.22), size=8.5, color=C_GRAY)
        tb(slide, f"📅 {period}", sx + Inches(0.64), by + Inches(0.52),
           sw - Inches(0.72), Inches(0.20), size=8.5, bold=True, color=pc)

    # ── 간트 차트 영역 ────────────────────────────────────
    gx, gy = Inches(0.18), Inches(2.08)
    gw, gh = W - Inches(0.36), Inches(3.22)
    section_card(slide, gx, gy, gw, gh, "📅  개발 일정 (2026년 4월 ~ 12월)", C_NAVY)

    # 월 헤더
    MONTHS = ["4월", "5월", "6월", "7월", "8월", "9월", "10월", "11월", "12월"]
    N = len(MONTHS)
    label_w  = Inches(2.6)
    chart_x  = gx + label_w + Inches(0.08)
    chart_w  = gw - label_w - Inches(0.2)
    mw       = chart_w / N
    mh_y     = gy + Inches(0.42)

    for mi, m in enumerate(MONTHS):
        mx = chart_x + mi * mw
        bg = RGBColor(0xD0, 0xD9, 0xF0) if mi < 3 else RGBColor(0xE2, 0xE8, 0xF0)
        rect(slide, mx + Inches(0.01), mh_y, mw - Inches(0.02), Inches(0.30), bg)
        tb(slide, m, mx + Inches(0.01), mh_y + Inches(0.01),
           mw - Inches(0.02), Inches(0.28),
           size=9, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)

    # 현재 시점 표시선 (5월 = index 1)
    cur_x = chart_x + 1.5 * mw   # 5월 중간
    rect(slide, cur_x, mh_y, Inches(0.02), gh - Inches(0.44),
         RGBColor(0xFF, 0x44, 0x44))
    tb(slide, "▲ 현재", cur_x - Inches(0.22), mh_y + gh - Inches(0.60),
       Inches(0.7), Inches(0.22), size=7.5, bold=True,
       color=RGBColor(0xFF, 0x44, 0x44), align=PP_ALIGN.CENTER)

    # 간트 행 데이터
    # start/end: 0=4월, 1=5월, ... 8=12월  (소수점 = 월 내 위치)
    gantt_rows = [
        # ── 1단계 ──
        {
            "group":  "1단계",
            "label":  "배차 효율화 개발",
            "sub":    "자동배차 기능 구현",
            "color":  C_GREEN,
            "bg":     RGBColor(0xE6, 0xF7, 0xF4),
            "start":  0.0,
            "end":    3.0,
            "pct":    70,
            "status": "진행중",
        },
        {
            "group":  "1단계",
            "label":  "  └ WMS 데이터 연동",
            "sub":    "출고문서·ERP 인터페이스",
            "color":  RGBColor(0x0F, 0x9D, 0x58),
            "bg":     RGBColor(0xF0, 0xFD, 0xF4),
            "start":  0.0,
            "end":    1.5,
            "pct":    100,
            "status": "완료",
        },
        {
            "group":  "1단계",
            "label":  "  └ 기준정보 관리",
            "sub":    "차량·납품처·운송경로",
            "color":  RGBColor(0x0F, 0x9D, 0x58),
            "bg":     RGBColor(0xF0, 0xFD, 0xF4),
            "start":  0.5,
            "end":    2.5,
            "pct":    80,
            "status": "진행중",
        },
        {
            "group":  "1단계",
            "label":  "  └ 자동배차 알고리즘",
            "sub":    "배차전략·최적화·분할",
            "color":  RGBColor(0x0F, 0x9D, 0x58),
            "bg":     RGBColor(0xF0, 0xFD, 0xF4),
            "start":  1.0,
            "end":    3.0,
            "pct":    50,
            "status": "진행중",
        },
        # ── 2단계 ──
        {
            "group":  "2단계",
            "label":  "차량관제·운송비 관리",
            "sub":    "관제·정산·리포트",
            "color":  C_ORANGE,
            "bg":     RGBColor(0xFF, 0xF0, 0xD9),
            "start":  3.0,
            "end":    9.0,
            "pct":    0,
            "status": "예정",
        },
        {
            "group":  "2단계",
            "label":  "  └ GPS 차량 위치 관제",
            "sub":    "실시간 추적·이상 알림",
            "color":  RGBColor(0xD9, 0x77, 0x06),
            "bg":     RGBColor(0xFE, 0xF9, 0xF0),
            "start":  3.0,
            "end":    5.5,
            "pct":    0,
            "status": "예정",
        },
        {
            "group":  "2단계",
            "label":  "  └ 운송비 자동 정산",
            "sub":    "경로·중량 기반 운임 계산",
            "color":  RGBColor(0xD9, 0x77, 0x06),
            "bg":     RGBColor(0xFE, 0xF9, 0xF0),
            "start":  4.5,
            "end":    7.5,
            "pct":    0,
            "status": "예정",
        },
        {
            "group":  "2단계",
            "label":  "  └ 실적·KPI 리포트",
            "sub":    "대시보드·분석 고도화",
            "color":  RGBColor(0xD9, 0x77, 0x06),
            "bg":     RGBColor(0xFE, 0xF9, 0xF0),
            "start":  7.0,
            "end":    9.0,
            "pct":    0,
            "status": "예정",
        },
    ]

    row_h = Inches(0.295)
    row_gap = Inches(0.04)
    for ri, row in enumerate(gantt_rows):
        ry2 = mh_y + Inches(0.32) + ri * (row_h + row_gap)

        is_main = not row["label"].startswith("  └")

        # 행 배경
        rect(slide, gx + Inches(0.1), ry2,
             gw - Inches(0.2), row_h,
             row["bg"] if is_main else RGBColor(0xF8, 0xFA, 0xFC))

        # 그룹 배지 (메인 행만)
        if is_main:
            gc = C_GREEN if row["group"] == "1단계" else C_ORANGE
            rect(slide, gx + Inches(0.12), ry2 + Inches(0.04),
                 Inches(0.52), Inches(0.20), gc)
            tb(slide, row["group"],
               gx + Inches(0.12), ry2 + Inches(0.04),
               Inches(0.52), Inches(0.20),
               size=7.5, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        # 레이블
        lx_off = Inches(0.70) if is_main else Inches(0.20)
        tb(slide, row["label"],
           gx + lx_off, ry2 + Inches(0.02),
           label_w - lx_off - Inches(0.05), Inches(0.16),
           size=8.5 if is_main else 8, bold=is_main, color=C_DARK)
        tb(slide, row["sub"],
           gx + lx_off, ry2 + Inches(0.17),
           label_w - lx_off - Inches(0.05), Inches(0.14),
           size=7.5, color=C_GRAY, italic=True)

        # 간트 바 배경
        rect(slide, chart_x, ry2 + Inches(0.06),
             chart_w, row_h - Inches(0.12),
             RGBColor(0xEE, 0xF2, 0xF7))

        # 계획 바
        bar_s = chart_x + row["start"] / N * chart_w
        bar_e = chart_x + row["end"]   / N * chart_w
        bar_l = bar_e - bar_s
        bar_h = row_h - Inches(0.12)
        bar_y = ry2 + Inches(0.06)

        rect(slide, bar_s, bar_y, bar_l, bar_h,
             RGBColor(0xBB, 0xD6, 0xBB) if row["group"] == "1단계"
             else RGBColor(0xF5, 0xD6, 0xA0))

        # 진행 바 (완료 부분)
        if row["pct"] > 0:
            prog_l = bar_l * row["pct"] / 100
            rect(slide, bar_s, bar_y, prog_l, bar_h, row["color"])
            if prog_l > Inches(0.3):
                tb(slide, f"{row['pct']}%",
                   bar_s + Inches(0.04), bar_y,
                   Inches(0.4), bar_h,
                   size=7.5, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

        # 상태 배지 (오른쪽 끝)
        sc_map = {"완료": C_GREEN, "진행중": C_ORANGE,
                  "예정": RGBColor(0x64, 0x74, 0x8B)}
        sc = sc_map.get(row["status"], C_GRAY)
        rect(slide, bar_e + Inches(0.08), bar_y + Inches(0.03),
             Inches(0.44), bar_h - Inches(0.06), sc)
        tb(slide, row["status"],
           bar_e + Inches(0.08), bar_y + Inches(0.03),
           Inches(0.44), bar_h - Inches(0.06),
           size=7, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # ── 하단: 단계별 주요 산출물 ─────────────────────────
    ox, oy = Inches(0.18), gy + gh + Inches(0.14)
    ow, oh = W - Inches(0.36), H - oy - Inches(0.14)
    section_card(slide, ox, oy, ow, oh, "📋  단계별 주요 산출물 및 목표", C_BLUE)

    stages = [
        {
            "title":  "1단계  |  배차 효율화 (4월 ~ 6월)",
            "color":  C_GREEN,
            "bg":     RGBColor(0xE8, 0xF8, 0xF4),
            "items":  [
                "✔ WMS·ERP 출고 데이터 연동 완료",
                "✔ 차량·납품처·운송경로 기준정보 관리",
                "✔ 고정차량 관리 (TMS_CARCLASS 16차급)",
                "✔ 자동배차 알고리즘 (용량·경로 최적화)",
                "✔ 배차전략 파라미터 설정 UI",
                "◷ 배차 결과 확정 및 지시서 출력",
            ]
        },
        {
            "title":  "2단계  |  차량관제·운송비 관리 (7월 ~ 12월)",
            "color":  C_ORANGE,
            "bg":     RGBColor(0xFF, 0xF3, 0xE0),
            "items":  [
                "◷ GPS 기반 실시간 차량 위치 관제",
                "◷ 이상 경보 및 배송 지연 알림",
                "◷ 경로·중량 기반 운송비 자동 정산",
                "◷ 차량별·노선별 운행 실적 분석",
                "◷ 월별·고객사별 KPI 대시보드",
                "◷ 운전자 모바일 배차지시 연동",
            ]
        },
    ]

    sw2 = (ow - Inches(0.32)) / 2
    for si, stg in enumerate(stages):
        sx2 = ox + Inches(0.14) + si * (sw2 + Inches(0.06))
        sy2 = oy + Inches(0.44)
        sh2 = oh - Inches(0.52)
        rect(slide, sx2, sy2, sw2, sh2, stg["bg"])
        rect(slide, sx2, sy2, sw2, Inches(0.26), stg["color"])
        tb(slide, stg["title"], sx2 + Inches(0.08), sy2 + Inches(0.02),
           sw2 - Inches(0.12), Inches(0.23),
           size=9.5, bold=True, color=C_WHITE)
        for ii, item in enumerate(stg["items"]):
            tb(slide, item,
               sx2 + Inches(0.12), sy2 + Inches(0.30) + ii * Inches(0.22),
               sw2 - Inches(0.18), Inches(0.21),
               size=9, color=C_DARK)

    # 범례
    lx2 = ox + Inches(0.14)
    ly2 = oy + oh - Inches(0.28)
    for color, label in [(C_GREEN, "완료"), (C_ORANGE, "진행중"),
                          (RGBColor(0x64, 0x74, 0x8B), "예정"),
                          (RGBColor(0xFF, 0x44, 0x44), "현재 시점")]:
        rect(slide, lx2, ly2 + Inches(0.06), Inches(0.18), Inches(0.13), color)
        tb(slide, label, lx2 + Inches(0.22), ly2,
           Inches(0.9), Inches(0.25), size=8, color=C_GRAY)
        lx2 += Inches(1.1)

    return slide


# ════════════════════════════════════════════════════════════════
#  메인
# ════════════════════════════════════════════════════════════════
def main():
    prs = new_prs()
    build_slide1(prs)
    build_slide2(prs)
    out = "/home/user/webapp/kleannara_TMS_개발진행현황_2026-05.pptx"
    prs.save(out)
    print(f"✅ 저장 완료: {out}")

if __name__ == "__main__":
    main()
