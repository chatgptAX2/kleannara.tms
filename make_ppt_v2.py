"""
TMS 개발 진행현황 보고자료 PPT v2 - 세련된 다크 테마 디자인
슬라이드1: TMS 시스템 소개 + 도입효과
슬라이드2: 개발 일정 (간트차트)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn, nsmap
from pptx.oxml import parse_xml
from lxml import etree
import copy

# ══════════════════════════════════════════════════════
#  디자인 토큰 (모던 다크 + 민트 포인트 테마)
# ══════════════════════════════════════════════════════
# 배경
BG          = RGBColor(0x0F, 0x17, 0x2A)   # 짙은 네이비 배경
BG2         = RGBColor(0x16, 0x21, 0x3E)   # 카드 배경
BG3         = RGBColor(0x1E, 0x2D, 0x50)   # 밝은 카드

# 포인트 컬러
MINT        = RGBColor(0x00, 0xD4, 0xAA)   # 민트 (1단계 / 완료)
AMBER       = RGBColor(0xFF, 0xB8, 0x00)   # 앰버 (2단계 / 진행)
SKY         = RGBColor(0x38, 0xBD, 0xF8)   # 스카이블루 (포인트)
LAVENDER    = RGBColor(0xA7, 0x8B, 0xFA)   # 라벤더 (효과카드3)
CORAL       = RGBColor(0xFF, 0x6B, 0x6B)   # 코랄 (현재선)
INDIGO      = RGBColor(0x60, 0x6C, 0xFC)   # 인디고 (효과카드1)

# 텍스트
TXT         = RGBColor(0xF1, 0xF5, 0xF9)   # 기본 텍스트
TXT2        = RGBColor(0x94, 0xA3, 0xB8)   # 보조 텍스트
TXT3        = RGBColor(0x64, 0x74, 0x8B)   # 흐린 텍스트
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

# 테두리/구분선
BORDER      = RGBColor(0x2D, 0x3D, 0x60)
BORDER2     = RGBColor(0x38, 0x4D, 0x72)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ══════════════════════════════════════════════════════
#  헬퍼 함수
# ══════════════════════════════════════════════════════
def new_prs():
    p = Presentation()
    p.slide_width  = SLIDE_W
    p.slide_height = SLIDE_H
    return p

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def R(slide, x, y, w, h, color, radius=0):
    """사각형 추가 (radius는 현재 미사용 — pptx 제한)"""
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def T(slide, text, x, y, w, h,
      size=11, bold=False, color=TXT,
      align=PP_ALIGN.LEFT, italic=False,
      wrap=True, line_spacing=None, space_before=0):
    """텍스트 박스"""
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    r = p.add_run()
    r.text = text
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return txb

def add_line(slide, x1, y1, x2, y2, color=BORDER, width_pt=0.5):
    """선 추가"""
    from pptx.util import Pt as UPt
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = UPt(width_pt)
    return connector

def pill_badge(slide, text, x, y, w, h, bg_color, txt_color=WHITE, size=7.5):
    """알약형 배지"""
    R(slide, x, y, w, h, bg_color)
    T(slide, text, x, y, w, h,
      size=size, bold=True, color=txt_color, align=PP_ALIGN.CENTER)

def dot_line(slide, x, y, length, color=BORDER):
    """수평 구분선"""
    R(slide, x, y, length, Inches(0.008), color)


# ══════════════════════════════════════════════════════
#  슬라이드 공통 헤더
# ══════════════════════════════════════════════════════
def common_header(slide, title_main, title_sub, page_str):
    W = SLIDE_W
    # 전체 배경
    R(slide, 0, 0, W, SLIDE_H, BG)

    # 상단 헤더 영역 (그라데이션 느낌: 두 레이어)
    R(slide, 0, 0, W, Inches(1.0), RGBColor(0x0A, 0x10, 0x1F))

    # 좌측 포인트 바 (4px 굵기)
    R(slide, 0, 0, Inches(0.055), SLIDE_H, MINT)

    # 상단 구분선
    R(slide, 0, Inches(1.0), W, Inches(0.006), MINT)

    # 브랜드 태그
    T(slide, "kleannara  ·  TMS",
      Inches(0.2), Inches(0.10), Inches(3.5), Inches(0.28),
      size=9, bold=False, color=TXT3, align=PP_ALIGN.LEFT)

    # 페이지 번호
    T(slide, page_str,
      Inches(12.8), Inches(0.10), Inches(0.42), Inches(0.28),
      size=9, color=TXT3, align=PP_ALIGN.RIGHT)

    # 메인 타이틀
    T(slide, title_main,
      Inches(0.2), Inches(0.25), Inches(10), Inches(0.52),
      size=24, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # 서브 타이틀
    T(slide, title_sub,
      Inches(0.2), Inches(0.72), Inches(8), Inches(0.26),
      size=9.5, color=MINT, align=PP_ALIGN.LEFT)


# ══════════════════════════════════════════════════════
#  슬라이드 1: TMS 시스템 소개 + 도입효과
# ══════════════════════════════════════════════════════
def build_slide1(prs):
    slide = blank(prs)
    common_header(slide,
                  "TMS (Transport Management System) 시스템 소개",
                  "물류 운송 자동화 · 배차 최적화 · 운송비 관리 통합 플랫폼",
                  "01")

    BODY_TOP = Inches(1.10)
    BODY_H   = SLIDE_H - BODY_TOP - Inches(0.12)

    # ── 좌측: TMS 정의 + 구성도 (폭 5.7) ─────────────────
    LX = Inches(0.18)
    LW = Inches(5.70)

    # ① 정의 카드
    cy = BODY_TOP + Inches(0.16)
    R(slide, LX, cy, LW, Inches(1.22), BG2)
    R(slide, LX, cy, Inches(0.04), Inches(1.22), SKY)   # 좌측 강조선

    T(slide, "TMS란?",
      LX + Inches(0.16), cy + Inches(0.08), LW - Inches(0.22), Inches(0.24),
      size=10, bold=True, color=SKY)
    T(slide,
      "Transport Management System — 물류 기업의 출하 계획부터 배차 지시,\n"
      "운송 실행, 운임 정산까지 전 과정을 통합·자동화하는 운송관리 시스템",
      LX + Inches(0.16), cy + Inches(0.32), LW - Inches(0.22), Inches(0.82),
      size=9.5, color=TXT, wrap=True, line_spacing=15)
    cy += Inches(1.34)

    # ② 시스템 연계 흐름도
    R(slide, LX, cy, LW, Inches(0.96), BG2)
    R(slide, LX, cy, Inches(0.04), Inches(0.96), AMBER)

    T(slide, "시스템 연계 구조",
      LX + Inches(0.16), cy + Inches(0.08), LW - Inches(0.22), Inches(0.22),
      size=9, bold=True, color=AMBER)

    nodes = [
        ("ERP",    "출고지시",  INDIGO),
        ("WMS",    "창고관리",  SKY),
        ("TMS",    "운송관리",  MINT),
        ("현장",   "배송실행",  AMBER),
    ]
    node_w = (LW - Inches(0.32)) / len(nodes)
    nx_start = LX + Inches(0.12)
    for ni, (name, sub, nc) in enumerate(nodes):
        nx = nx_start + ni * node_w
        # 노드 박스
        R(slide, nx + Inches(0.06), cy + Inches(0.34),
          node_w - Inches(0.16), Inches(0.46), BG3)
        R(slide, nx + Inches(0.06), cy + Inches(0.34),
          node_w - Inches(0.16), Inches(0.06), nc)
        T(slide, name,
          nx + Inches(0.06), cy + Inches(0.40),
          node_w - Inches(0.16), Inches(0.22),
          size=10, bold=True, color=nc, align=PP_ALIGN.CENTER)
        T(slide, sub,
          nx + Inches(0.06), cy + Inches(0.60),
          node_w - Inches(0.16), Inches(0.18),
          size=7.5, color=TXT2, align=PP_ALIGN.CENTER)
        # 화살표
        if ni < len(nodes) - 1:
            T(slide, "›",
              nx + node_w - Inches(0.14), cy + Inches(0.44),
              Inches(0.22), Inches(0.28),
              size=14, bold=True, color=TXT3, align=PP_ALIGN.CENTER)
    cy += Inches(1.08)

    # ③ 핵심 기능 4개
    features = [
        (MINT,     "배차 자동화",
         "출고 데이터 기반 차량 용량·경로·우선순위를\n자동 계산, 최적 배차안 즉시 생성"),
        (SKY,      "차량 통합 관리",
         "차량 제원(톤수·파렛트·적재중량)·차급 기준\n등록 및 배차 자동 매핑"),
        (AMBER,    "운송경로·비용",
         "납품처별 배송구역·우선순위·경유 경로 등록\n운송경로 비용 기준 운임 자동 산출"),
        (LAVENDER, "실적 분석",
         "배차 실적·노선별 운임·차량 가동률 KPI를\n대시보드에서 실시간 모니터링"),
    ]

    fh = (BODY_H - (cy - BODY_TOP) - Inches(0.18)) / len(features) - Inches(0.06)
    for fi, (fc, ftitle, fdesc) in enumerate(features):
        fy = cy + fi * (fh + Inches(0.06))
        R(slide, LX, fy, LW, fh, BG2)
        # 좌측 컬러 바
        R(slide, LX, fy, Inches(0.04), fh, fc)
        # 아이콘 원형 배경
        R(slide, LX + Inches(0.14), fy + Inches(0.10),
          Inches(0.30), fh - Inches(0.20), BG3)
        # 번호
        T(slide, f"0{fi+1}",
          LX + Inches(0.14), fy + Inches(0.10),
          Inches(0.30), fh - Inches(0.20),
          size=9, bold=True, color=fc, align=PP_ALIGN.CENTER)
        T(slide, ftitle,
          LX + Inches(0.54), fy + Inches(0.08),
          LW - Inches(0.62), Inches(0.24),
          size=9.5, bold=True, color=WHITE)
        T(slide, fdesc,
          LX + Inches(0.54), fy + Inches(0.28),
          LW - Inches(0.62), fh - Inches(0.32),
          size=8.5, color=TXT2, wrap=True, line_spacing=13)

    # ── 우측: 도입 기대효과 (폭 7.3) ──────────────────────
    RX = LX + LW + Inches(0.20)
    RW = SLIDE_W - RX - Inches(0.18)

    # 섹션 레이블
    T(slide, "도입 기대효과",
      RX, BODY_TOP + Inches(0.16), RW, Inches(0.28),
      size=10, bold=True, color=MINT)
    dot_line(slide, RX, BODY_TOP + Inches(0.44), RW, MINT)

    effects = [
        {
            "num":    "01",
            "title":  "배차 업무 효율화",
            "color":  INDIGO,
            "bg":     RGBColor(0x1A, 0x20, 0x4A),
            "kpi":    "업무시간 70% 절감",
            "points": [
                "수작업 배차 → 자동배차 전환으로 담당자 업무 부하 대폭 감소",
                "배차 처리 시간 대폭 단축 및 배차 품질 표준화",
                "오배차·누락 사전 차단, 배차 이력 통합 관리",
            ],
        },
        {
            "num":    "02",
            "title":  "운송비 절감",
            "color":  MINT,
            "bg":     RGBColor(0x0D, 0x2A, 0x26),
            "kpi":    "운송비 15~25% 절감",
            "points": [
                "혼적·적재율 최적화로 공차율 감소 및 차량 운용 효율 향상",
                "경로 최적화로 불필요한 이동 거리·연료비 절감",
                "운송경로별 비용 기준 자동 정산으로 운임 투명성 확보",
            ],
        },
        {
            "num":    "03",
            "title":  "고객 서비스 향상",
            "color":  SKY,
            "bg":     RGBColor(0x0D, 0x23, 0x35),
            "kpi":    "납기 준수율 향상",
            "points": [
                "납품처별 배송 우선순위 반영으로 납기 약속 준수율 향상",
                "배송 현황 실시간 조회 및 이상 발생 시 즉각 대응",
                "정확한 배송 예정 시간(ETA) 안내로 고객 신뢰도 제고",
            ],
        },
        {
            "num":    "04",
            "title":  "경영 의사결정 지원",
            "color":  AMBER,
            "bg":     RGBColor(0x2A, 0x1F, 0x07),
            "kpi":    "물류 전체 가시성 확보",
            "points": [
                "배차 실적·운임·차량 가동률 KPI 대시보드 실시간 제공",
                "노선별·차급별 비용 분석으로 운송 전략 수립 지원",
                "ERP·WMS 데이터 통합 연동, 물류 전체 가시성 확보",
            ],
        },
    ]

    eff_top = BODY_TOP + Inches(0.52)
    eff_total_h = BODY_H - Inches(0.52) - Inches(0.06)
    eff_h = eff_total_h / len(effects) - Inches(0.07)

    for ei, eff in enumerate(effects):
        ey = eff_top + ei * (eff_h + Inches(0.07))
        # 카드 배경
        R(slide, RX, ey, RW, eff_h, eff["bg"])
        # 상단 컬러 라인
        R(slide, RX, ey, RW, Inches(0.04), eff["color"])

        # 번호 배지
        R(slide, RX + Inches(0.14), ey + Inches(0.12),
          Inches(0.36), Inches(0.36), eff["color"])
        T(slide, eff["num"],
          RX + Inches(0.14), ey + Inches(0.12),
          Inches(0.36), Inches(0.36),
          size=10, bold=True, color=BG, align=PP_ALIGN.CENTER)

        # 제목
        T(slide, eff["title"],
          RX + Inches(0.60), ey + Inches(0.10),
          Inches(2.4), Inches(0.26),
          size=11, bold=True, color=WHITE)

        # KPI 뱃지
        kpi_w = Inches(1.9)
        R(slide, RX + RW - kpi_w - Inches(0.14), ey + Inches(0.10),
          kpi_w, Inches(0.26),
          RGBColor(
              int(eff["color"][0] * 0.3),
              int(eff["color"][1] * 0.3),
              int(eff["color"][2] * 0.3),
          ))
        T(slide, f"▶  {eff['kpi']}",
          RX + RW - kpi_w - Inches(0.14), ey + Inches(0.10),
          kpi_w, Inches(0.26),
          size=8, bold=True, color=eff["color"], align=PP_ALIGN.CENTER)

        # 구분선
        R(slide, RX + Inches(0.60), ey + Inches(0.38),
          RW - Inches(0.72), Inches(0.005), BORDER)

        # 포인트 텍스트
        for pi, pt in enumerate(eff["points"]):
            T(slide, f"·  {pt}",
              RX + Inches(0.60), ey + Inches(0.42) + pi * Inches(0.22),
              RW - Inches(0.72), Inches(0.21),
              size=8.8, color=TXT2, wrap=False)

    return slide


# ══════════════════════════════════════════════════════
#  슬라이드 2: 개발 일정
# ══════════════════════════════════════════════════════
def build_slide2(prs):
    slide = blank(prs)
    common_header(slide,
                  "TMS 개발 일정 및 단계별 추진 계획",
                  "1단계: 배차 효율화 (4~6월)  ·  2단계: 차량관제 및 운송비 관리 (7~12월)",
                  "02")

    W, H = SLIDE_W, SLIDE_H
    BODY_TOP = Inches(1.10)

    # ── 단계 요약 카드 2개 ───────────────────────────────
    SW = (W - Inches(0.36) - Inches(0.16)) / 2
    for si, (phase, title, desc, period, pc, bg) in enumerate([
        ("1단계", "배차 효율화 개발",
         "자동배차 기능을 통한 배차 효율화 향상",
         "2026. 04 ~ 06월", MINT,
         RGBColor(0x0D, 0x2A, 0x22)),
        ("2단계", "차량관제 및 운송비 관리 개발",
         "실시간 차량관제 · 운임 정산 · 실적 분석",
         "2026. 07 ~ 12월", AMBER,
         RGBColor(0x2A, 0x1E, 0x07)),
    ]):
        sx = Inches(0.18) + si * (SW + Inches(0.16))
        sy = BODY_TOP + Inches(0.12)
        sh = Inches(0.88)
        R(slide, sx, sy, SW, sh, bg)
        R(slide, sx, sy, SW, Inches(0.05), pc)
        # 단계 배지
        R(slide, sx + Inches(0.16), sy + Inches(0.16),
          Inches(0.58), Inches(0.56), pc)
        T(slide, phase,
          sx + Inches(0.16), sy + Inches(0.16),
          Inches(0.58), Inches(0.56),
          size=10, bold=True, color=BG, align=PP_ALIGN.CENTER)
        T(slide, title,
          sx + Inches(0.86), sy + Inches(0.10),
          SW - Inches(1.05), Inches(0.28),
          size=11, bold=True, color=WHITE)
        T(slide, desc,
          sx + Inches(0.86), sy + Inches(0.36),
          SW - Inches(1.05), Inches(0.22),
          size=8.8, color=TXT2)
        T(slide, f"📅  {period}",
          sx + Inches(0.86), sy + Inches(0.58),
          SW - Inches(1.05), Inches(0.22),
          size=9, bold=True, color=pc)

    # ── 간트차트 영역 ────────────────────────────────────
    GX = Inches(0.18)
    GY = BODY_TOP + Inches(1.14)
    GW = W - Inches(0.36)
    GH = Inches(3.52)

    R(slide, GX, GY, GW, GH, BG2)
    R(slide, GX, GY, GW, Inches(0.04), MINT)

    # 레이블
    T(slide, "개발 일정",
      GX + Inches(0.18), GY + Inches(0.10), Inches(2), Inches(0.26),
      size=10, bold=True, color=MINT)

    # 월 헤더 설정
    MONTHS    = ["4월","5월","6월","7월","8월","9월","10월","11월","12월"]
    N         = len(MONTHS)
    LABEL_W   = Inches(2.86)
    CHART_X   = GX + LABEL_W
    CHART_W   = GW - LABEL_W - Inches(0.16)
    MW        = CHART_W / N
    HDR_Y     = GY + Inches(0.44)
    HDR_H     = Inches(0.32)

    # 월 헤더 셀
    for mi, m in enumerate(MONTHS):
        mx = CHART_X + mi * MW
        is_phase1 = mi < 3
        hbg = RGBColor(0x0D, 0x2A, 0x22) if is_phase1 else RGBColor(0x2A, 0x1F, 0x07)
        hfg = MINT if is_phase1 else AMBER
        R(slide, mx + Inches(0.01), HDR_Y,
          MW - Inches(0.02), HDR_H, hbg)
        T(slide, m,
          mx + Inches(0.01), HDR_Y,
          MW - Inches(0.02), HDR_H,
          size=9, bold=True, color=hfg, align=PP_ALIGN.CENTER)

    # 현재 시점선 (5월 중반)
    CUR_X = CHART_X + 1.55 * MW
    R(slide, CUR_X - Inches(0.008), HDR_Y,
      Inches(0.016), GH - Inches(0.46), CORAL)
    T(slide, "TODAY",
      CUR_X - Inches(0.32), GY + GH - Inches(0.32),
      Inches(0.64), Inches(0.22),
      size=7, bold=True, color=CORAL, align=PP_ALIGN.CENTER)

    # 간트 행 데이터
    gantt_rows = [
        # ─── 1단계 메인 ───
        {
            "label":   "1단계  배차 효율화 개발",
            "sub":     "자동배차 기능을 통한 배차 효율화 향상",
            "color":   MINT,
            "bg_row":  RGBColor(0x0D, 0x22, 0x1C),
            "start":   0.0, "end": 3.0,
            "pct":     70,
            "status":  "진행중",
            "is_main": True,
        },
        {
            "label":   "  WMS · ERP 출고 데이터 연동",
            "sub":     "출고문서(SHPDH/SHPDI) · ERP 인터페이스(IFWMS113)",
            "color":   RGBColor(0x34, 0xD3, 0x99),
            "bg_row":  BG2,
            "start":   0.0, "end": 1.5,
            "pct":     100,
            "status":  "완료",
            "is_main": False,
        },
        {
            "label":   "  기준정보 관리",
            "sub":     "차량·납품처·운송경로·고정차량(TMS_CARCLASS 16차급)",
            "color":   RGBColor(0x34, 0xD3, 0x99),
            "bg_row":  BG2,
            "start":   0.3, "end": 2.6,
            "pct":     85,
            "status":  "진행중",
            "is_main": False,
        },
        {
            "label":   "  자동배차 알고리즘",
            "sub":     "배차전략 설정 · 용량/경로 최적화 · 배차 분할",
            "color":   RGBColor(0x34, 0xD3, 0x99),
            "bg_row":  BG2,
            "start":   1.0, "end": 3.0,
            "pct":     55,
            "status":  "진행중",
            "is_main": False,
        },
        # ─── 2단계 메인 ───
        {
            "label":   "2단계  차량관제 및 운송비 관리",
            "sub":     "실시간 차량관제 · 운임 정산 · 실적 분석",
            "color":   AMBER,
            "bg_row":  RGBColor(0x22, 0x18, 0x06),
            "start":   3.0, "end": 9.0,
            "pct":     0,
            "status":  "예정",
            "is_main": True,
        },
        {
            "label":   "  GPS 실시간 차량 위치 관제",
            "sub":     "실시간 위치 추적 · 이상 경보 · 배송 지연 알림",
            "color":   RGBColor(0xFF, 0xC9, 0x4D),
            "bg_row":  BG2,
            "start":   3.0, "end": 5.8,
            "pct":     0,
            "status":  "예정",
            "is_main": False,
        },
        {
            "label":   "  운송비 자동 정산",
            "sub":     "경로·중량 기반 운임 자동 계산 · 비용 분석",
            "color":   RGBColor(0xFF, 0xC9, 0x4D),
            "bg_row":  BG2,
            "start":   4.8, "end": 7.6,
            "pct":     0,
            "status":  "예정",
            "is_main": False,
        },
        {
            "label":   "  실적 KPI 대시보드",
            "sub":     "월별·노선별·차량별 KPI · 대시보드 고도화",
            "color":   RGBColor(0xFF, 0xC9, 0x4D),
            "bg_row":  BG2,
            "start":   7.2, "end": 9.0,
            "pct":     0,
            "status":  "예정",
            "is_main": False,
        },
    ]

    ROW_H   = Inches(0.348)
    ROW_GAP = Inches(0.024)
    ROWS_Y  = HDR_Y + HDR_H + Inches(0.06)

    for ri, row in enumerate(gantt_rows):
        ry = ROWS_Y + ri * (ROW_H + ROW_GAP)
        is_main = row["is_main"]

        # 행 배경
        R(slide, GX + Inches(0.04), ry,
          GW - Inches(0.08), ROW_H, row["bg_row"])

        # 좌측 컬러 선 (메인 행)
        if is_main:
            R(slide, GX + Inches(0.04), ry,
              Inches(0.025), ROW_H, row["color"])

        # 레이블 텍스트
        label_x = GX + Inches(0.16) if is_main else GX + Inches(0.30)
        T(slide, row["label"],
          label_x, ry + Inches(0.04),
          LABEL_W - Inches(0.38), Inches(0.20),
          size=9 if is_main else 8.5,
          bold=is_main, color=WHITE if is_main else TXT)
        T(slide, row["sub"],
          label_x, ry + Inches(0.22),
          LABEL_W - Inches(0.38), Inches(0.14),
          size=7.5, color=TXT3, italic=True)

        # 차트 영역 배경 (격자 스트라이프)
        for mi in range(N):
            stripe_c = RGBColor(0x14, 0x1E, 0x38) if mi % 2 == 0 else RGBColor(0x12, 0x1A, 0x32)
            R(slide, CHART_X + mi * MW, ry,
              MW, ROW_H, stripe_c)

        # 간트 바
        bar_s = CHART_X + row["start"] / N * CHART_W
        bar_e = CHART_X + row["end"]   / N * CHART_W
        bar_l = bar_e - bar_s
        bar_y = ry + Inches(0.07)
        bar_h = ROW_H - Inches(0.14)

        # 계획 바 (연한 배경)
        plan_c = (RGBColor(0x0A, 0x3D, 0x30) if row["color"] == MINT or
                  row["color"] == RGBColor(0x34, 0xD3, 0x99)
                  else RGBColor(0x3D, 0x2A, 0x06))
        R(slide, bar_s, bar_y, bar_l, bar_h, plan_c)

        # 진행 바
        if row["pct"] > 0:
            prog = bar_l * row["pct"] / 100
            R(slide, bar_s, bar_y, prog, bar_h, row["color"])
            if prog > Inches(0.35):
                T(slide, f"{row['pct']}%",
                  bar_s + Inches(0.06), bar_y,
                  Inches(0.4), bar_h,
                  size=7.5, bold=True, color=BG, align=PP_ALIGN.LEFT)

        # 상태 배지
        st_colors = {
            "완료":  (MINT,                          BG),
            "진행중": (AMBER,                         BG),
            "예정":  (RGBColor(0x3D, 0x4D, 0x6B),   TXT2),
        }
        sc, stc = st_colors.get(row["status"], (TXT3, BG))
        badge_x = bar_e + Inches(0.08)
        R(slide, badge_x, bar_y + Inches(0.02),
          Inches(0.50), bar_h - Inches(0.04), sc)
        T(slide, row["status"],
          badge_x, bar_y + Inches(0.02),
          Inches(0.50), bar_h - Inches(0.04),
          size=7.5, bold=True, color=stc, align=PP_ALIGN.CENTER)

    # ── 하단: 주요 산출물 ────────────────────────────────
    OX = Inches(0.18)
    OY = GY + GH + Inches(0.14)
    OW = W - Inches(0.36)
    OH = H - OY - Inches(0.12)

    R(slide, OX, OY, OW, OH, BG2)
    R(slide, OX, OY, OW, Inches(0.04), SKY)

    T(slide, "단계별 주요 산출물",
      OX + Inches(0.18), OY + Inches(0.10), Inches(2.4), Inches(0.24),
      size=9.5, bold=True, color=SKY)

    delivs = [
        {
            "phase": "1단계  |  4 ~ 6월",
            "color": MINT,
            "items": [
                ("✔", "WMS·ERP 출고 데이터 연동"),
                ("✔", "차량·납품처·운송경로 기준정보"),
                ("✔", "고정차량 관리 (16차급)"),
                ("✔", "자동배차 알고리즘"),
                ("◷", "배차 결과 확정·지시서 출력"),
            ]
        },
        {
            "phase": "2단계  |  7 ~ 12월",
            "color": AMBER,
            "items": [
                ("◷", "GPS 실시간 차량 위치 관제"),
                ("◷", "이상경보·배송 지연 알림"),
                ("◷", "경로·중량 기반 운임 자동 정산"),
                ("◷", "차량별·노선별 운행 실적 분석"),
                ("◷", "KPI 대시보드 고도화"),
            ]
        },
    ]

    col_w = (OW - Inches(0.36)) / 2
    for di, deliv in enumerate(delivs):
        dx = OX + Inches(0.14) + di * (col_w + Inches(0.10))
        dy = OY + Inches(0.38)

        T(slide, deliv["phase"],
          dx, dy, col_w, Inches(0.22),
          size=9, bold=True, color=deliv["color"])

        for ii, (mark, item_text) in enumerate(deliv["items"]):
            item_y = dy + Inches(0.26) + ii * Inches(0.20)
            mc = MINT if mark == "✔" else TXT3
            T(slide, mark,
              dx, item_y, Inches(0.22), Inches(0.20),
              size=9, bold=True, color=mc, align=PP_ALIGN.CENTER)
            T(slide, item_text,
              dx + Inches(0.24), item_y, col_w - Inches(0.26), Inches(0.20),
              size=8.8, color=TXT)

    # 수직 구분선
    mid = OX + Inches(0.14) + col_w + Inches(0.05)
    R(slide, mid, OY + Inches(0.34), Inches(0.008), OH - Inches(0.42), BORDER)

    # 범례 (우측 하단)
    lex = OX + OW - Inches(4.0)
    ley = OY + OH - Inches(0.24)
    for lc, lt in [(MINT, "완료"), (AMBER, "진행중"),
                   (RGBColor(0x3D, 0x4D, 0x6B), "예정"),
                   (CORAL, "현재 시점")]:
        R(slide, lex, ley + Inches(0.05), Inches(0.16), Inches(0.14), lc)
        T(slide, lt, lex + Inches(0.20), ley, Inches(0.75), Inches(0.22),
          size=7.5, color=TXT3)
        lex += Inches(0.98)

    return slide


# ══════════════════════════════════════════════════════
#  실행
# ══════════════════════════════════════════════════════
def main():
    prs = new_prs()
    build_slide1(prs)
    build_slide2(prs)
    out = "/home/user/webapp/kleannara_TMS_개발진행현황_v2_2026-05.pptx"
    prs.save(out)
    print(f"✅  {out}")

if __name__ == "__main__":
    main()
